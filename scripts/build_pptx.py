"""
JSON → 편집 가능한 PPTX 생성기.
- 슬라이드 캔버스: 13.333 × 7.5 inch (16:9, 1920×1080 매핑)
- 디자인 토큰은 design_system/tokens.json 단일 소스에서 가져옴
- 타입별 슬라이드 빌더: cover / chapter_divider / content / card_grid_4 / stage_flow / alert_close
"""
from __future__ import annotations
import json
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


ROOT = Path(__file__).resolve().parent.parent
DESIGN = ROOT / "design_system"
OUTPUTS = ROOT / "outputs"

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
EDGE_IN    = 0.7    # 100 px @ 1920 ≈ 0.7 inch
TOP_IN     = 0.65

FONT_FAMILY = "Pretendard"
FONT_FALLBACK = "Pretendard Variable"
FONT_SAFE = "Apple SD Gothic Neo"  # macOS fallback


def _resolve_speaker_photo(sp: dict) -> str:
    """speaker.photo 없으면 이름으로 assets/<name>*.{jpg,png} 자동 매칭 (HTML과 동일 로직)."""
    if sp.get("photo"):
        return sp["photo"]
    name = (sp.get("name") or "").strip()
    if not name:
        return ""
    asset_dir = ROOT / "assets"
    if not asset_dir.exists():
        return ""
    variants = [f"{name} 프로필", name, f"{name}_profile", f"{name}_프로필"]
    exts = ["jpg", "JPG", "jpeg", "JPEG", "png", "PNG"]
    for variant in variants:
        for ext in exts:
            cand = asset_dir / f"{variant}.{ext}"
            if cand.exists():
                return f"assets/{cand.name}"
    return ""


def _set_font_typeface(run, name: str = FONT_FAMILY) -> None:
    """Run의 latin + ea(한글) + cs 폰트를 모두 설정.
    python-pptx의 run.font.name 은 latin만 설정하기 때문에
    한글이 East Asian 기본 폰트로 떨어지는 문제가 있음 → 직접 OOXML 조작."""
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        node = rPr.find(qn(tag))
        if node is None:
            node = etree.SubElement(rPr, qn(tag))
        node.set("typeface", name)


# ---------- Token loading ----------

with open(DESIGN / "tokens.json", encoding="utf-8") as f:
    TOKENS = json.load(f)
COLORS = TOKENS["colors"]
TYPE   = TOKENS["type"]


def hex_to_rgb(h: str) -> RGBColor:
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def C(name: str) -> RGBColor:
    return hex_to_rgb(COLORS[name])


# ---------- Inline markup parser ----------

INLINE_RE = re.compile(r"(<em>.*?</em>|<b>.*?</b>)", re.DOTALL)
TAG_RE    = re.compile(r"<(em|b)>(.*?)</\1>", re.DOTALL)


def parse_inline(text: str):
    """Split text into runs: [(text, {"em":bool,"b":bool}), ...]
    Supports <em> (accent color) and <b> (bold). Newlines preserved."""
    if text is None:
        return [("", {})]
    parts = INLINE_RE.split(text)
    runs = []
    for p in parts:
        if not p:
            continue
        m = TAG_RE.match(p)
        if m:
            tag, inner = m.group(1), m.group(2)
            attrs = {"em": tag == "em", "b": tag == "b"}
            runs.append((inner, attrs))
        else:
            runs.append((p, {}))
    return runs


# ---------- Low-level drawing helpers ----------

def add_text(slide, x, y, w, h, text, *, size_pt, weight=400, color="ink",
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, leading=1.3,
             tracking_em=0, allow_inline=True):
    """Add a text box with given style. Supports <em>/<b> inline markup."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    # Handle multi-line: split by \n
    lines = str(text).split("\n")
    for li, line in enumerate(lines):
        p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = leading
        # Clear any existing runs from default empty paragraph
        for old_r in list(p.runs):
            old_r._r.getparent().remove(old_r._r)
        runs = parse_inline(line) if allow_inline else [(line, {})]
        for j, (txt, attrs) in enumerate(runs):
            r = p.add_run()
            r.text = txt
            _set_font_typeface(r, FONT_FAMILY)  # latin + ea(한글) + cs 모두
            r.font.size = Pt(size_pt)
            r.font.bold = (weight >= 700) or attrs.get("b", False)
            run_color = "accent" if attrs.get("em") else color
            r.font.color.rgb = C(run_color)
            # tracking via XML <a:rPr spc="-N"> (1/100 pt)
            if tracking_em:
                spc = int(tracking_em * size_pt * 100)
                rPr = r._r.get_or_add_rPr()
                rPr.set("spc", str(spc))
    return tb


def add_rect(slide, x, y, w, h, *, fill="paper", line=None, line_w_pt=0.75):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = C(fill)
    if line:
        shp.line.color.rgb = C(line)
        shp.line.width = Pt(line_w_pt)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_hline(slide, x, y, w, *, color="line", weight_pt=0.75):
    ln = slide.shapes.add_connector(1, Inches(x), Inches(y), Inches(x + w), Inches(y))
    ln.line.color.rgb = C(color)
    ln.line.width = Pt(weight_pt)
    return ln


def add_footer(slide, meta, page_no, total):
    talk = meta.get("talk_title_in_footer", "")
    add_hline(slide, EDGE_IN, SLIDE_H_IN - 0.55, SLIDE_W_IN - 2 * EDGE_IN, color="line")
    add_text(slide, EDGE_IN, SLIDE_H_IN - 0.5, 9, 0.3,
             "ZEROBASE", size_pt=10, weight=700, color="ink",
             tracking_em=-0.01, allow_inline=False)
    add_text(slide, EDGE_IN + 1.0, SLIDE_H_IN - 0.5, 9, 0.3,
             f"·  {talk}", size_pt=10, weight=400, color="muted",
             allow_inline=False)
    add_text(slide, SLIDE_W_IN - EDGE_IN - 1.5, SLIDE_H_IN - 0.5, 1.5, 0.3,
             f"{page_no:02d} / {total:02d}", size_pt=10, weight=500,
             color="muted", align=PP_ALIGN.RIGHT, allow_inline=False)


def add_chapter_tag(slide, x, y, text, small=False):
    sz = 14 if small else 18
    return add_text(slide, x, y, SLIDE_W_IN - 2 * EDGE_IN, 0.45,
                    text, size_pt=sz, weight=700, color="accent",
                    tracking_em=-0.01, allow_inline=False)


# ---------- Slide builders ----------

def slide_cover(prs, meta):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    add_text(s, EDGE_IN + 0.2, 1.4, SLIDE_W_IN - 2 * EDGE_IN, 0.5,
             meta.get("event_label", ""),
             size_pt=14, weight=700, color="accent",
             tracking_em=0.18, allow_inline=False)
    title_text = meta.get("title", "")
    em_text = meta.get("title_em", "")
    add_text(s, EDGE_IN + 0.2, 2.1, SLIDE_W_IN - 2 * EDGE_IN, 1.3,
             title_text, size_pt=56, weight=800, color="ink",
             tracking_em=-0.02, leading=1.05, allow_inline=False)
    if em_text:
        add_text(s, EDGE_IN + 0.2, 3.3, SLIDE_W_IN - 2 * EDGE_IN, 1.3,
                 em_text, size_pt=56, weight=800, color="accent",
                 tracking_em=-0.02, leading=1.05, allow_inline=False)
    add_text(s, EDGE_IN + 0.2, 4.7, SLIDE_W_IN - 2 * EDGE_IN, 1.4,
             meta.get("subtitle", ""), size_pt=20, weight=500,
             color="muted_2", leading=1.4, allow_inline=False)
    # bottom meta band
    add_hline(s, EDGE_IN + 0.2, SLIDE_H_IN - 1.2, SLIDE_W_IN - 2 * (EDGE_IN + 0.2),
              color="line")
    sp = meta.get("speaker", {}) or {}
    add_text(s, EDGE_IN + 0.2, SLIDE_H_IN - 1.05, 7, 0.4,
             sp.get("name", ""), size_pt=18, weight=700, color="ink",
             allow_inline=False)
    add_text(s, EDGE_IN + 0.2, SLIDE_H_IN - 0.65, 9, 0.35,
             sp.get("title", ""), size_pt=12, weight=500, color="muted_2",
             allow_inline=False)
    add_text(s, SLIDE_W_IN - EDGE_IN - 4.2, SLIDE_H_IN - 1.05, 4, 0.4,
             f"{meta.get('date','')}  ·  {meta.get('session','')}",
             size_pt=14, weight=500, color="muted_2", align=PP_ALIGN.RIGHT,
             allow_inline=False)


def slide_chapter_divider(prs, meta, s_data, page_no, total):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    bg = add_rect(s, 0, 0, SLIDE_W_IN, SLIDE_H_IN, fill="paper_2")
    bg.shadow.inherit = False
    add_text(s, EDGE_IN + 0.4, 2.1, SLIDE_W_IN - 2 * EDGE_IN, 0.5,
             s_data.get("ch", ""), size_pt=18, weight=700,
             color="accent", tracking_em=0.2, allow_inline=False)
    # 타이틀 축소 (72 → 60pt) + max-width로 한 줄 유도
    add_text(s, EDGE_IN + 0.4, 2.8, SLIDE_W_IN - 2 * EDGE_IN, 2.5,
             s_data.get("title", ""), size_pt=60, weight=800,
             color="ink", tracking_em=-0.02, leading=1.1,
             allow_inline=False)
    # sub: title 줄 수에 따라 동적 위치 (1줄→4.15 / 2줄→5.4 / 3줄→6.0)
    if s_data.get("sub"):
        title_text = s_data.get("title", "")
        n_lines = title_text.count("\n") + 1
        sub_y = 4.15 if n_lines == 1 else 5.4 if n_lines == 2 else 6.0
        add_text(s, EDGE_IN + 0.4, sub_y, SLIDE_W_IN - 2 * EDGE_IN, 1.2,
                 s_data["sub"], size_pt=24, weight=500, color="muted_2",
                 leading=1.45, allow_inline=False)
    add_footer(s, meta, page_no, total)


def slide_content(prs, meta, s_data, page_no, total):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    add_chapter_tag(s, EDGE_IN, TOP_IN, s_data.get("chapter", ""), small=True)
    add_text(s, EDGE_IN, TOP_IN + 0.4, SLIDE_W_IN - 2 * EDGE_IN, 1.6,
             s_data.get("title", ""), size_pt=40, weight=800,
             color="ink", tracking_em=-0.02, leading=1.08)
    add_text(s, EDGE_IN, TOP_IN + 1.7, SLIDE_W_IN - 2 * EDGE_IN, 0.9,
             s_data.get("lede", ""), size_pt=18, weight=500,
             color="muted_2", leading=1.45, allow_inline=False)
    # bullets
    bullets = s_data.get("bullets") or []
    by = TOP_IN + 3.1
    for b in bullets:
        # bullet mark
        bx = EDGE_IN + 0.05
        bm = add_rect(s, bx, by + 0.18, 0.18, 0.03, fill="accent")
        bm.line.fill.background()
        add_text(s, bx + 0.32, by, SLIDE_W_IN - 2 * EDGE_IN - 0.4, 0.6,
                 b, size_pt=18, weight=400, color="ink", leading=1.5,
                 allow_inline=False)
        by += 0.55
    add_footer(s, meta, page_no, total)


def slide_card_grid(prs, meta, s_data, page_no, total):
    """N-up pastel cards (2~5)."""
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    add_chapter_tag(s, EDGE_IN, TOP_IN, s_data.get("chapter", ""), small=True)
    title_h = 1.4
    add_text(s, EDGE_IN, TOP_IN + 0.4, SLIDE_W_IN - 2 * EDGE_IN, title_h,
             s_data.get("title", ""), size_pt=36, weight=800,
             color="ink", tracking_em=-0.02, leading=1.08)
    lede_y = TOP_IN + 0.4 + title_h
    if s_data.get("lede"):
        add_text(s, EDGE_IN, lede_y, SLIDE_W_IN - 2 * EDGE_IN, 0.7,
                 s_data["lede"], size_pt=16, weight=500, color="muted_2",
                 leading=1.4, allow_inline=False)
        cy = lede_y + 0.95
    else:
        cy = lede_y + 0.45
    cards = (s_data.get("cards") or [])[:7]
    n = max(1, min(7, len(cards)))
    inner_gap = 0.12 if n >= 6 else 0.18
    available = SLIDE_W_IN - 2 * EDGE_IN
    card_w = (available - inner_gap * (n - 1)) / n
    card_h = SLIDE_H_IN - cy - 0.85
    # Adjust font sizes for narrow cards
    title_pt = 13 if n >= 6 else 18 if n >= 4 else 22 if n == 3 else 24
    item_pt  = 10 if n >= 6 else 12 if n >= 4 else 13 if n == 3 else 14
    item_gap = 0.32 if n >= 6 else 0.38 if n >= 4 else 0.42
    for i, c in enumerate(cards):
        cx = EDGE_IN + i * (card_w + inner_gap)
        bg = add_rect(s, cx, cy, card_w, card_h, fill="pastel_blue")
        bg.line.fill.background()
        add_text(s, cx + 0.25, cy + 0.25, card_w - 0.5, 0.35,
                 c.get("num", ""), size_pt=11, weight=700, color="accent",
                 tracking_em=0.16, allow_inline=False)
        add_text(s, cx + 0.25, cy + 0.6, card_w - 0.5, 1.0,
                 c.get("title", ""), size_pt=title_pt, weight=700, color="ink",
                 leading=1.25, allow_inline=False)
        add_hline(s, cx + 0.25, cy + 1.55, card_w - 0.5, color="line", weight_pt=0.75)
        iy = cy + 1.7
        for it in (c.get("items") or [])[:6]:
            mark = add_rect(s, cx + 0.28, iy + 0.12, 0.13, 0.025, fill="accent")
            mark.line.fill.background()
            add_text(s, cx + 0.5, iy, card_w - 0.7, item_gap + 0.1,
                     it, size_pt=item_pt, weight=400, color="muted_2",
                     leading=1.4, allow_inline=False)
            iy += item_gap
    add_footer(s, meta, page_no, total)


# alias
slide_card_grid_4 = slide_card_grid


def slide_two_col(prs, meta, s_data, page_no, total):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    add_chapter_tag(s, EDGE_IN, TOP_IN, s_data.get("chapter", ""), small=True)
    add_text(s, EDGE_IN, TOP_IN + 0.4, SLIDE_W_IN - 2 * EDGE_IN, 1.4,
             s_data.get("title", ""), size_pt=36, weight=800,
             color="ink", tracking_em=-0.02, leading=1.08)
    lede_y = TOP_IN + 1.8
    if s_data.get("lede"):
        add_text(s, EDGE_IN, lede_y, SLIDE_W_IN - 2 * EDGE_IN, 0.7,
                 s_data["lede"], size_pt=16, weight=500, color="muted_2",
                 leading=1.4, allow_inline=False)
        cy = lede_y + 0.95
    else:
        cy = lede_y + 0.4
    cols = (s_data.get("cols") or [])[:2]
    if not cols:
        return
    gap = 0.32
    available = SLIDE_W_IN - 2 * EDGE_IN
    col_w = (available - gap) / 2
    col_h = SLIDE_H_IN - cy - 0.85
    for i, col in enumerate(cols):
        cx = EDGE_IN + i * (col_w + gap)
        is_accent = bool(col.get("accent"))
        bg = add_rect(s, cx, cy, col_w, col_h,
                      fill=("pastel_blue" if is_accent else "paper"),
                      line=(None if is_accent else "line"),
                      line_w_pt=0.75)
        if is_accent:
            bg.line.fill.background()
        ix = cx + 0.32
        iy = cy + 0.32
        if col.get("num"):
            add_text(s, ix, iy, col_w - 0.64, 0.35,
                     col["num"], size_pt=11, weight=700, color="accent",
                     tracking_em=0.16, allow_inline=False)
            iy += 0.45
        add_text(s, ix, iy, col_w - 0.64, 0.95,
                 col.get("title", ""), size_pt=22, weight=700, color="ink",
                 leading=1.25, allow_inline=False)
        iy += 0.95
        add_hline(s, ix, iy + 0.05, col_w - 0.64, color="line", weight_pt=0.75)
        iy += 0.25
        if col.get("lede"):
            add_text(s, ix, iy, col_w - 0.64, 0.85,
                     col["lede"], size_pt=15, weight=500, color="muted_2",
                     leading=1.5, allow_inline=False)
            iy += 0.85
        for it in (col.get("items") or [])[:8]:
            mark = add_rect(s, ix, iy + 0.13, 0.13, 0.025, fill="accent")
            mark.line.fill.background()
            add_text(s, ix + 0.22, iy, col_w - 0.86, 0.55,
                     it, size_pt=15, weight=400, color="ink",
                     leading=1.45, allow_inline=False)
            iy += 0.45
    add_footer(s, meta, page_no, total)


def slide_toc(prs, meta, s_data, page_no, total):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    add_chapter_tag(s, EDGE_IN, TOP_IN, s_data.get("chapter", "Contents"), small=True)
    add_text(s, EDGE_IN, TOP_IN + 0.4, SLIDE_W_IN - 2 * EDGE_IN, 1.6,
             s_data.get("title", "목차"), size_pt=44, weight=800,
             color="ink", tracking_em=-0.02, leading=1.08)
    items = s_data.get("items") or []
    cy = TOP_IN + 2.4
    avail = SLIDE_W_IN - 2 * EDGE_IN
    col_w = (avail - 0.5) / 2
    row_h = 0.62
    rows_per_col = (len(items) + 1) // 2
    for i, it in enumerate(items):
        col_idx = i // rows_per_col
        row_idx = i % rows_per_col
        x = EDGE_IN + col_idx * (col_w + 0.5)
        y = cy + row_idx * row_h
        add_text(s, x, y, 0.7, row_h - 0.1,
                 f"{i+1:02d}", size_pt=18, weight=700, color="accent",
                 allow_inline=False)
        add_text(s, x + 0.7, y, col_w - 0.7, row_h - 0.1,
                 it, size_pt=20, weight=500, color="ink",
                 tracking_em=-0.01, leading=1.3, allow_inline=False)
        add_hline(s, x, y + row_h - 0.05, col_w, color="line", weight_pt=0.5)
    add_footer(s, meta, page_no, total)


def slide_speaker(prs, meta, s_data, page_no, total):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    sp = s_data.get("speaker") or meta.get("speaker", {}) or {}
    add_chapter_tag(s, EDGE_IN, TOP_IN, s_data.get("chapter", "강사 소개"), small=True)
    add_text(s, EDGE_IN, TOP_IN + 0.4, SLIDE_W_IN - 2 * EDGE_IN, 1.4,
             s_data.get("title", "오늘의 연사"), size_pt=40, weight=800,
             color="ink", tracking_em=-0.02, leading=1.08)
    cy = TOP_IN + 2.2
    # 사진 영역: 3:4 portrait 비율 (1.95 x 2.6 inch)
    photo_w = 1.95
    photo_h = 2.6
    photo_path = sp.get("photo") or _resolve_speaker_photo(sp)
    if photo_path:
        # 박스/프레임 제거 — 사진만 깔끔하게 (cover 방식: 비율 유지하며 영역 가득 채움, 가운데 위쪽 정렬)
        img_p = _resolve_image(photo_path)
        if img_p and img_p.exists():
            _add_picture_cover(s, img_p, EDGE_IN, cy, photo_w, photo_h,
                               focus_y_ratio=0.2)  # 인물 머리쪽 살짝 위
        else:
            # missing image fallback
            bg = add_rect(s, EDGE_IN, cy, photo_w, photo_h, fill="paper_2", line="line")
            add_text(s, EDGE_IN, cy, photo_w, photo_h,
                     "?", size_pt=72, weight=800, color="muted",
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, allow_inline=False)
    else:
        # 사진 없을 때: 이름 첫 글자 fallback (박스 유지)
        bg = add_rect(s, EDGE_IN, cy, photo_w, photo_h, fill="paper_2", line="line")
        add_text(s, EDGE_IN, cy, photo_w, photo_h,
                 sp.get("name", "?")[:1], size_pt=72, weight=800, color="muted",
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, allow_inline=False)
    info_x = EDGE_IN + photo_w + 0.4
    info_w = SLIDE_W_IN - info_x - EDGE_IN
    add_text(s, info_x, cy, info_w, 0.9,
             sp.get("name", ""), size_pt=40, weight=800, color="ink",
             tracking_em=-0.02, leading=1.05, allow_inline=False)
    add_text(s, info_x, cy + 0.85, info_w, 0.55,
             sp.get("title", ""), size_pt=18, weight=700, color="accent",
             tracking_em=-0.01, allow_inline=False)
    # tags
    iy = cy + 1.5
    tx = info_x
    for t in sp.get("tags", []):
        # Approximate width by character count
        tw = max(0.7, len(t) * 0.15 + 0.4)
        if tx + tw > info_x + info_w:
            tx = info_x; iy += 0.4
        bg2 = add_rect(s, tx, iy, tw, 0.32, fill="accent_soft")
        bg2.line.fill.background()
        add_text(s, tx, iy, tw, 0.32, t, size_pt=12, weight=600,
                 color="accent", align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 allow_inline=False)
        tx += tw + 0.08
    iy += 0.55
    for b in sp.get("bio", [])[:6]:
        mark = add_rect(s, info_x, iy + 0.13, 0.13, 0.025, fill="accent")
        mark.line.fill.background()
        add_text(s, info_x + 0.22, iy, info_w - 0.22, 0.55,
                 b, size_pt=14, weight=400, color="ink",
                 leading=1.5, allow_inline=False)
        iy += 0.42
    add_footer(s, meta, page_no, total)


def slide_qa_close(prs, meta, s_data, page_no, total):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    has_contact = bool(s_data.get("contact"))
    has_thanks  = bool(s_data.get("thanks"))
    # Lift Q&A up if there's contact/thanks below
    qa_y = SLIDE_H_IN/2 - 3.2 if (has_contact or has_thanks) else SLIDE_H_IN/2 - 2.5
    add_text(s, 0, qa_y, SLIDE_W_IN, 3.0,
             "Q&A", size_pt=160, weight=800, color="accent",
             tracking_em=-0.04, leading=1.0,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             allow_inline=False)
    cur_y = qa_y + 3.0
    if s_data.get("sub"):
        add_text(s, 0, cur_y + 0.05, SLIDE_W_IN, 0.7,
                 s_data["sub"], size_pt=20, weight=500, color="muted_2",
                 align=PP_ALIGN.CENTER, allow_inline=False)
        cur_y += 0.7
    if has_contact:
        # Centered "boxed" contact
        contact = s_data["contact"]
        # rough width estimate
        w = max(2.5, len(contact) * 0.18 + 0.8)
        x = (SLIDE_W_IN - w) / 2
        bg = add_rect(s, x, cur_y + 0.25, w, 0.6, fill="paper", line="line")
        add_text(s, x, cur_y + 0.25, w, 0.6, contact,
                 size_pt=18, weight=700, color="ink",
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 tracking_em=-0.01, allow_inline=False)
        cur_y += 0.95
    if has_thanks:
        add_text(s, 0, cur_y + 0.2, SLIDE_W_IN, 0.6,
                 s_data["thanks"], size_pt=16, weight=500, color="muted_2",
                 align=PP_ALIGN.CENTER, allow_inline=False)
    add_footer(s, meta, page_no, total)


def slide_compare_table(prs, meta, s_data, page_no, total):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    add_chapter_tag(s, EDGE_IN, TOP_IN, s_data.get("chapter", ""), small=True)
    add_text(s, EDGE_IN, TOP_IN + 0.4, SLIDE_W_IN - 2 * EDGE_IN, 1.4,
             s_data.get("title", ""), size_pt=36, weight=800,
             color="ink", tracking_em=-0.02, leading=1.08)
    headers = s_data.get("headers") or ["구분", "A", "B"]
    rows    = s_data.get("rows") or []
    accent_col = s_data.get("accent_col", -1)
    if s_data.get("lede"):
        add_text(s, EDGE_IN, TOP_IN + 1.6, SLIDE_W_IN - 2 * EDGE_IN, 0.7,
                 s_data["lede"], size_pt=16, weight=500, color="muted_2",
                 leading=1.4, allow_inline=False)
        cy = TOP_IN + 2.55
    else:
        cy = TOP_IN + 2.2
    avail_w = SLIDE_W_IN - 2 * EDGE_IN
    label_w = 1.95
    cell_w  = (avail_w - label_w) / 2
    head_h = 0.55
    row_h  = 0.7  # adjust if many rows
    n_rows = len(rows)
    concl  = s_data.get("conclusion")
    bottom_pad = 1.0 if concl else 0.85
    avail_h = SLIDE_H_IN - cy - bottom_pad
    table_h = avail_h - (0.3 if concl else 0)
    if n_rows > 0:
        row_h = (table_h - head_h) / n_rows
        row_h = min(row_h, 1.0)
    # Outer border
    outer = add_rect(s, EDGE_IN, cy, avail_w, head_h + row_h * n_rows,
                     fill="paper", line="line", line_w_pt=0.75)
    # Header row
    add_rect(s, EDGE_IN, cy, label_w, head_h, fill="paper_2", line="line", line_w_pt=0.5)
    add_text(s, EDGE_IN, cy, label_w, head_h, headers[0],
             size_pt=11, weight=700, color="accent",
             tracking_em=0.16, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             allow_inline=False)
    add_rect(s, EDGE_IN + label_w, cy, cell_w, head_h,
             fill="paper_2", line="line", line_w_pt=0.5)
    add_text(s, EDGE_IN + label_w, cy, cell_w, head_h, headers[1],
             size_pt=14, weight=700, color="ink",
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, allow_inline=False)
    add_rect(s, EDGE_IN + label_w + cell_w, cy, cell_w, head_h,
             fill="paper_2", line="line", line_w_pt=0.5)
    add_text(s, EDGE_IN + label_w + cell_w, cy, cell_w, head_h, headers[2],
             size_pt=14, weight=700, color="ink",
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, allow_inline=False)
    # Body rows
    ry = cy + head_h
    for r in rows:
        cells = r if isinstance(r, list) else [r.get("label",""), r.get("a",""), r.get("b","")]
        # label col
        add_rect(s, EDGE_IN, ry, label_w, row_h,
                 fill="paper_2", line="line", line_w_pt=0.5)
        add_text(s, EDGE_IN + 0.2, ry, label_w - 0.4, row_h,
                 cells[0], size_pt=11, weight=700, color="accent",
                 tracking_em=0.12, anchor=MSO_ANCHOR.MIDDLE, allow_inline=False)
        # A
        a_fill = "pastel_blue" if accent_col == 1 else "paper"
        add_rect(s, EDGE_IN + label_w, ry, cell_w, row_h,
                 fill=a_fill, line="line", line_w_pt=0.5)
        add_text(s, EDGE_IN + label_w + 0.25, ry, cell_w - 0.5, row_h,
                 cells[1], size_pt=14, weight=500, color="ink",
                 leading=1.4, anchor=MSO_ANCHOR.MIDDLE, allow_inline=False)
        # B
        b_fill = "pastel_blue" if accent_col == 2 else "paper"
        add_rect(s, EDGE_IN + label_w + cell_w, ry, cell_w, row_h,
                 fill=b_fill, line="line", line_w_pt=0.5)
        add_text(s, EDGE_IN + label_w + cell_w + 0.25, ry, cell_w - 0.5, row_h,
                 cells[2], size_pt=14, weight=500, color="ink",
                 leading=1.4, anchor=MSO_ANCHOR.MIDDLE, allow_inline=False)
        ry += row_h
    # Conclusion bar
    if concl:
        cb_y = ry + 0.18
        cb_h = 0.5
        add_rect(s, EDGE_IN, cb_y, avail_w, cb_h, fill="ink_2", line=None)
        add_text(s, EDGE_IN + 0.3, cb_y, avail_w - 0.6, cb_h, concl,
                 size_pt=14, weight=600, color="paper",
                 anchor=MSO_ANCHOR.MIDDLE, leading=1.4)
    add_footer(s, meta, page_no, total)


def slide_stage_flow(prs, meta, s_data, page_no, total):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    add_chapter_tag(s, EDGE_IN, TOP_IN, s_data.get("chapter", ""), small=True)
    add_text(s, EDGE_IN, TOP_IN + 0.4, SLIDE_W_IN - 2 * EDGE_IN, 1.6,
             s_data.get("title", ""), size_pt=40, weight=800,
             color="ink", tracking_em=-0.02, leading=1.08)
    if s_data.get("lede"):
        add_text(s, EDGE_IN, TOP_IN + 1.85, SLIDE_W_IN - 2 * EDGE_IN, 0.7,
                 s_data["lede"], size_pt=16, weight=500, color="muted_2",
                 leading=1.4, allow_inline=False)
    steps = s_data.get("steps") or []
    hi = s_data.get("highlight_index", -1)
    n = max(1, len(steps))
    arrow_w = 0.45
    inner_gap = 0.0
    available = SLIDE_W_IN - 2 * EDGE_IN
    box_w = (available - arrow_w * (n - 1)) / n
    box_h = 2.2
    by = TOP_IN + (3.3 if not s_data.get("lede") else 3.6)
    for i, st in enumerate(steps):
        x = EDGE_IN + i * (box_w + arrow_w)
        is_hi = (i == hi)
        bg = add_rect(s, x, by, box_w, box_h,
                      fill=("ink_2" if is_hi else "paper"),
                      line=(None if is_hi else "line"),
                      line_w_pt=0.75)
        if is_hi:
            bg.line.fill.background()
        add_text(s, x + 0.3, by + 0.3, box_w - 0.6, 0.4,
                 st.get("num", ""), size_pt=11, weight=700,
                 color=("paper" if is_hi else "accent"),
                 tracking_em=0.18, allow_inline=False)
        add_text(s, x + 0.3, by + 0.7, box_w - 0.6, 0.7,
                 st.get("name", ""), size_pt=22, weight=700,
                 color=("paper" if is_hi else "ink"),
                 tracking_em=-0.01, allow_inline=False)
        add_text(s, x + 0.3, by + 1.45, box_w - 0.6, 0.7,
                 st.get("desc", ""), size_pt=13, weight=400,
                 color=("paper" if is_hi else "muted_2"),
                 leading=1.4, allow_inline=False)
        if i < len(steps) - 1:
            ax = x + box_w
            add_text(s, ax, by + box_h / 2 - 0.3, arrow_w, 0.6,
                     "→", size_pt=22, weight=400, color="line_2",
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                     allow_inline=False)
    add_footer(s, meta, page_no, total)


def _resolve_image(rel_or_abs: str) -> Path:
    """프로젝트 루트 기준 상대경로 → 절대경로."""
    if not rel_or_abs:
        return Path()
    p = Path(rel_or_abs)
    if p.is_absolute() and p.exists():
        return p
    return ROOT / rel_or_abs


def _add_picture_cover(slide, img_path: Path, x_in, y_in, w_in, h_in, focus_y_ratio=0.5):
    """CSS object-fit: cover 동등 — 비율 유지하며 영역 가득, 넘치는 부분은 자름.
    focus_y_ratio: 세로 길이가 더 클 때 위/아래 어디를 살릴지 (0=상단, 1=하단, 0.5=가운데)."""
    if not img_path or not img_path.exists():
        ph = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(x_in), Inches(y_in),
                                    Inches(w_in), Inches(h_in))
        ph.fill.solid(); ph.fill.fore_color.rgb = C("paper_2")
        ph.line.color.rgb = C("line"); ph.line.width = Pt(0.75)
        return None
    from PIL import Image, ImageOps
    from io import BytesIO
    with Image.open(img_path) as im:
        im = ImageOps.exif_transpose(im)  # phone shots: auto-rotate
        if im.mode != "RGB":
            im = im.convert("RGB")
        iw, ih = im.size
        target_ratio = w_in / h_in
        img_ratio = iw / ih
        if img_ratio > target_ratio:
            # 이미지가 더 넓음 → 좌우 자름 (가운데 기준)
            new_w = int(ih * target_ratio)
            offset_x = (iw - new_w) // 2
            cropped = im.crop((offset_x, 0, offset_x + new_w, ih))
        else:
            # 이미지가 더 김 → 상하 자름 (focus_y_ratio 기준)
            new_h = int(iw / target_ratio)
            offset_y = max(0, min(ih - new_h, int((ih - new_h) * focus_y_ratio)))
            cropped = im.crop((0, offset_y, iw, offset_y + new_h))
        buf = BytesIO()
        cropped.save(buf, "JPEG", quality=92)
        buf.seek(0)
        return slide.shapes.add_picture(buf, Inches(x_in), Inches(y_in),
                                         width=Inches(w_in), height=Inches(h_in))


def _add_picture_fit(slide, img_path: Path, x_in, y_in, max_w_in, max_h_in):
    """add_picture 후 가로/세로 비율 보존하여 박스에 맞추고 가운데 정렬."""
    if not img_path or not img_path.exists():
        # placeholder rect
        ph = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(x_in), Inches(y_in),
                                    Inches(max_w_in), Inches(max_h_in))
        ph.fill.solid(); ph.fill.fore_color.rgb = C("paper_2")
        ph.line.color.rgb = C("line"); ph.line.width = Pt(0.75)
        slide.shapes.add_textbox(Inches(x_in), Inches(y_in),
                                 Inches(max_w_in), Inches(max_h_in)).text_frame.text = \
            f"(이미지 없음: {img_path})"
        return None
    # First try fit-to-width
    pic = slide.shapes.add_picture(str(img_path), Inches(x_in), Inches(y_in),
                                   width=Inches(max_w_in))
    if pic.height > Inches(max_h_in):
        sp = pic._element
        sp.getparent().remove(sp)
        pic = slide.shapes.add_picture(str(img_path), Inches(x_in), Inches(y_in),
                                       height=Inches(max_h_in))
    # Center within bounding box
    actual_w_in = pic.width / 914400
    actual_h_in = pic.height / 914400
    pic.left = Inches(x_in + (max_w_in - actual_w_in) / 2)
    pic.top  = Inches(y_in + (max_h_in - actual_h_in) / 2)
    return pic


def slide_image_full(prs, meta, s_data, page_no, total):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    # background
    bg_color = "paper_2" if s_data.get("bg") == "light" else "ink"
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                            Inches(SLIDE_W_IN), Inches(SLIDE_H_IN))
    bg.fill.solid(); bg.fill.fore_color.rgb = C(bg_color)
    bg.line.fill.background()
    bg.shadow.inherit = False

    has_caption = bool(s_data.get("caption"))
    caption_h = 0.85 if has_caption else 0
    img_top = 0.4
    img_bottom = SLIDE_H_IN - caption_h - 0.4
    img_max_h = img_bottom - img_top
    img_max_w = SLIDE_W_IN - 0.8

    img_path = _resolve_image(s_data.get("image", ""))
    _add_picture_fit(s, img_path, 0.4, img_top, img_max_w, img_max_h)

    if has_caption:
        cy = SLIDE_H_IN - caption_h
        cap_bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(cy),
                                    Inches(SLIDE_W_IN), Inches(caption_h))
        cap_bg.fill.solid(); cap_bg.fill.fore_color.rgb = C("paper")
        cap_bg.line.fill.background()
        add_text(s, EDGE_IN, cy + 0.18, SLIDE_W_IN - 2 * EDGE_IN, caption_h - 0.3,
                 s_data.get("caption", ""), size_pt=14, weight=500, color="muted_2",
                 leading=1.45)
    add_footer(s, meta, page_no, total)


def slide_image_headed(prs, meta, s_data, page_no, total):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    add_chapter_tag(s, EDGE_IN, TOP_IN, s_data.get("chapter", ""), small=True)
    add_text(s, EDGE_IN, TOP_IN + 0.4, SLIDE_W_IN - 2 * EDGE_IN, 1.4,
             s_data.get("title", ""), size_pt=36, weight=800,
             color="ink", tracking_em=-0.02, leading=1.08)
    cy = TOP_IN + 1.9
    if s_data.get("lede"):
        add_text(s, EDGE_IN, cy, SLIDE_W_IN - 2 * EDGE_IN, 0.7,
                 s_data["lede"], size_pt=16, weight=500, color="muted_2",
                 leading=1.4, allow_inline=False)
        cy += 0.85
    img_box_y = cy
    img_box_h = SLIDE_H_IN - img_box_y - 0.85
    img_box_w = SLIDE_W_IN - 2 * EDGE_IN
    # Bordered frame
    frame = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(EDGE_IN), Inches(img_box_y),
                               Inches(img_box_w), Inches(img_box_h))
    frame.fill.solid(); frame.fill.fore_color.rgb = C("paper_2")
    frame.line.color.rgb = C("line"); frame.line.width = Pt(0.75)
    img_path = _resolve_image(s_data.get("image", ""))
    _add_picture_fit(s, img_path,
                     EDGE_IN + 0.25, img_box_y + 0.25,
                     img_box_w - 0.5, img_box_h - 0.5)
    add_footer(s, meta, page_no, total)


def slide_image_split(prs, meta, s_data, page_no, total):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    add_chapter_tag(s, EDGE_IN, TOP_IN, s_data.get("chapter", ""), small=True)
    add_text(s, EDGE_IN, TOP_IN + 0.4, SLIDE_W_IN - 2 * EDGE_IN, 1.4,
             s_data.get("title", ""), size_pt=36, weight=800,
             color="ink", tracking_em=-0.02, leading=1.08)
    cy = TOP_IN + 2.2
    avail_w = SLIDE_W_IN - 2 * EDGE_IN
    avail_h = SLIDE_H_IN - cy - 0.85
    gap = 0.32
    pane_w = (avail_w - gap) / 2
    pos = s_data.get("image_position", "left")
    img_x = EDGE_IN if pos == "left" else EDGE_IN + pane_w + gap
    txt_x = EDGE_IN + pane_w + gap if pos == "left" else EDGE_IN
    # image frame
    frame = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(img_x), Inches(cy),
                               Inches(pane_w), Inches(avail_h))
    frame.fill.solid(); frame.fill.fore_color.rgb = C("paper_2")
    frame.line.color.rgb = C("line"); frame.line.width = Pt(0.75)
    img_path = _resolve_image(s_data.get("image", ""))
    _add_picture_fit(s, img_path,
                     img_x + 0.2, cy + 0.2,
                     pane_w - 0.4, avail_h - 0.4)
    # text pane
    ty = cy + 0.2
    if s_data.get("body_title"):
        add_text(s, txt_x, ty, pane_w, 1.0,
                 s_data["body_title"], size_pt=22, weight=700, color="ink",
                 leading=1.25)
        ty += 0.95
    if s_data.get("body_lede"):
        add_text(s, txt_x, ty, pane_w, 1.0,
                 s_data["body_lede"], size_pt=15, weight=500, color="muted_2",
                 leading=1.5, allow_inline=False)
        ty += 0.9
    for it in (s_data.get("items") or [])[:8]:
        mark = add_rect(s, txt_x, ty + 0.13, 0.13, 0.025, fill="accent")
        mark.line.fill.background()
        add_text(s, txt_x + 0.22, ty, pane_w - 0.22, 0.55,
                 it, size_pt=15, weight=400, color="ink",
                 leading=1.45, allow_inline=False)
        ty += 0.45
    add_footer(s, meta, page_no, total)


def slide_priority_matrix(prs, meta, s_data, page_no, total):
    """기호(◎/○/△) 매트릭스 + 행 강조 + 결론 배너 (PPTX)."""
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    add_chapter_tag(s, EDGE_IN, TOP_IN, s_data.get("chapter", ""), small=True)
    add_text(s, EDGE_IN, TOP_IN + 0.4, SLIDE_W_IN - 2 * EDGE_IN, 1.2,
             s_data.get("title", ""), size_pt=36, weight=800,
             color="ink", tracking_em=-0.02, leading=1.08)
    has_lede = bool(s_data.get("lede"))
    if has_lede:
        add_text(s, EDGE_IN, TOP_IN + 1.55, SLIDE_W_IN - 2 * EDGE_IN, 1.0,
                 s_data.get("lede", ""), size_pt=15, weight=500,
                 color="muted_2", leading=1.45, allow_inline=True)
        cy = TOP_IN + 2.4
    else:
        cy = TOP_IN + 2.0

    headers = s_data.get("headers", [])
    rows = s_data.get("rows", [])
    n_cols = len(headers)
    highlight_row = s_data.get("highlight_row", -1)
    avail_w = SLIDE_W_IN - 2 * EDGE_IN
    avail_h = SLIDE_H_IN - cy - 1.6
    label_w = 1.7
    cell_w = (avail_w - label_w) / max(1, n_cols - 1)
    head_h = 0.55
    row_h = min(0.7, (avail_h - head_h - 0.9) / max(1, len(rows)))

    # legend (top-right, above matrix)
    legend = "◎ 핵심   ○ 보조   △ 옵션"
    add_text(s, EDGE_IN + avail_w - 3.0, cy - 0.35, 3.0, 0.3,
             legend, size_pt=11, weight=500, color="muted_2",
             align=PP_ALIGN.RIGHT, allow_inline=False)

    # header row
    for k, h in enumerate(headers):
        x = EDGE_IN + (0 if k == 0 else label_w + (k - 1) * cell_w)
        w = label_w if k == 0 else cell_w
        rect = add_rect(s, x, cy, w, head_h, fill="ink_2")
        add_text(s, x, cy, w, head_h, h, size_pt=12, weight=700,
                 color="paper", align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 allow_inline=False)
    # rows
    for r_idx, row in enumerate(rows):
        is_hl = r_idx == highlight_row
        ry = cy + head_h + r_idx * row_h
        for k, cell in enumerate(row):
            x = EDGE_IN + (0 if k == 0 else label_w + (k - 1) * cell_w)
            w = label_w if k == 0 else cell_w
            fill = "paper_2" if k == 0 else "paper"
            line_color = "line"
            rect = add_rect(s, x, ry, w, row_h, fill=fill, line=line_color, line_w_pt=0.5)
            if is_hl:
                rect.line.color.rgb = C("ink_2")
                rect.line.width = Pt(1.5)
                # PPTX dashed: set dash style
                from pptx.enum.dml import MSO_LINE_DASH_STYLE
                try:
                    rect.line.dash_style = MSO_LINE_DASH_STYLE.DASH
                except Exception:
                    pass
            text = str(cell)
            font_size = 13 if k == 0 else 16
            weight = 700 if k == 0 else 600
            color = "ink" if not (k > 0 and text in ("△",)) else "muted"
            add_text(s, x, ry, w, row_h, text, size_pt=font_size, weight=weight,
                     color=color, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                     allow_inline=False)

    # conclusion banner
    conclusion = s_data.get("conclusion", "")
    if conclusion:
        cb_y = cy + head_h + len(rows) * row_h + 0.35
        cb_w = avail_w * 0.72
        cb_x = EDGE_IN + (avail_w - cb_w) / 2
        cb_h = 0.6
        rect = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(cb_x), Inches(cb_y),
                                  Inches(cb_w), Inches(cb_h))
        rect.adjustments[0] = 0.5
        rect.fill.background()
        rect.line.color.rgb = C("ink_2")
        rect.line.width = Pt(2)
        rect.shadow.inherit = False
        add_text(s, cb_x, cb_y, cb_w, cb_h, conclusion,
                 size_pt=14, weight=600, color="accent",
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, allow_inline=True)
    add_footer(s, meta, page_no, total)


def slide_step_compare(prs, meta, s_data, page_no, total):
    """좌우 카드 + 숫자 뱃지 단계 비교 (PPTX)."""
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    add_chapter_tag(s, EDGE_IN, TOP_IN, s_data.get("chapter", ""), small=True)
    add_text(s, EDGE_IN, TOP_IN + 0.4, SLIDE_W_IN - 2 * EDGE_IN, 1.2,
             s_data.get("title", ""), size_pt=36, weight=800,
             color="ink", tracking_em=-0.02, leading=1.08)
    if s_data.get("lede"):
        add_text(s, EDGE_IN, TOP_IN + 1.55, SLIDE_W_IN - 2 * EDGE_IN, 1.0,
                 s_data.get("lede", ""), size_pt=15, weight=500,
                 color="muted_2", leading=1.45, allow_inline=True)
        cy = TOP_IN + 2.4
    else:
        cy = TOP_IN + 2.0

    cols = s_data.get("cols", [])[:2]
    avail_w = SLIDE_W_IN - 2 * EDGE_IN
    gap = 0.32
    col_w = (avail_w - gap) / 2
    has_footer = bool(s_data.get("footer"))
    avail_h = SLIDE_H_IN - cy - (1.4 if has_footer else 1.0)
    for j, col in enumerate(cols):
        x = EDGE_IN + j * (col_w + gap)
        bg = "ink_2" if col.get("accent") else "ink"
        rect = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(x), Inches(cy),
                                  Inches(col_w), Inches(avail_h))
        rect.adjustments[0] = 0.06
        rect.fill.solid()
        rect.fill.fore_color.rgb = C(bg)
        rect.line.fill.background()
        rect.shadow.inherit = False
        # head
        pad = 0.32
        head_y = cy + pad
        add_text(s, x + pad, head_y, col_w - 2*pad, 0.4,
                 col.get("label", ""), size_pt=15, weight=800,
                 color="paper", allow_inline=False)
        if col.get("sub"):
            add_text(s, x + pad, head_y + 0.42, col_w - 2*pad, 0.3,
                     col.get("sub", ""), size_pt=11, weight=400,
                     color="paper", allow_inline=False)
            list_y = head_y + 0.85
        else:
            list_y = head_y + 0.55
        # divider line
        add_hline(s, x + pad, list_y - 0.12, col_w - 2*pad, color="paper", weight_pt=0.5)
        # items
        items = col.get("items", [])
        item_h = (avail_h - (list_y - cy) - pad) / max(1, len(items))
        for k, it in enumerate(items):
            iy = list_y + k * item_h
            num = it.get("num", f"{k+1:02d}")
            # num badge
            badge_size = 0.36
            badge = add_rect(s, x + pad, iy + 0.05, badge_size, badge_size,
                             fill="paper")
            badge.fill.fore_color.rgb = C("paper")
            badge.fill.transparency = 0  # solid
            try:
                from pptx.dml.color import RGBColor
                # alpha approach not directly supported, use lighter shade
                badge.fill.fore_color.rgb = C("ink_3")
            except Exception:
                pass
            add_text(s, x + pad, iy + 0.05, badge_size, badge_size, num,
                     size_pt=10, weight=700, color="paper",
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                     allow_inline=False)
            # title + desc
            add_text(s, x + pad + badge_size + 0.18, iy, col_w - pad*2 - badge_size - 0.18, 0.42,
                     it.get("title", ""), size_pt=14, weight=700, color="paper",
                     allow_inline=True)
            if it.get("desc"):
                add_text(s, x + pad + badge_size + 0.18, iy + 0.42, col_w - pad*2 - badge_size - 0.18, 0.4,
                         it.get("desc", ""), size_pt=11, weight=400, color="paper",
                         allow_inline=True)
    # footer
    if has_footer:
        fy = cy + avail_h + 0.25
        add_text(s, EDGE_IN, fy, avail_w, 0.5,
                 s_data.get("footer", ""), size_pt=14, weight=600, color="accent",
                 align=PP_ALIGN.CENTER, allow_inline=True)
    add_footer(s, meta, page_no, total)


def slide_before_after(prs, meta, s_data, page_no, total):
    """Before/After + 가운데 화살표 (PPTX)."""
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    add_chapter_tag(s, EDGE_IN, TOP_IN, s_data.get("chapter", ""), small=True)
    add_text(s, EDGE_IN, TOP_IN + 0.4, SLIDE_W_IN - 2 * EDGE_IN, 1.2,
             s_data.get("title", ""), size_pt=36, weight=800,
             color="ink", tracking_em=-0.02, leading=1.08)
    if s_data.get("lede"):
        add_text(s, EDGE_IN, TOP_IN + 1.55, SLIDE_W_IN - 2 * EDGE_IN, 1.0,
                 s_data.get("lede", ""), size_pt=15, weight=500,
                 color="muted_2", leading=1.45, allow_inline=True)
        cy = TOP_IN + 2.4
    else:
        cy = TOP_IN + 2.0

    avail_w = SLIDE_W_IN - 2 * EDGE_IN
    arrow_w = 0.55
    box_w = (avail_w - arrow_w) / 2
    box_h = 2.4
    note = s_data.get("note", "")
    source = s_data.get("source", "")

    # before box
    before = s_data.get("before", {})
    bx = EDGE_IN
    rect = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(bx), Inches(cy),
                              Inches(box_w), Inches(box_h))
    rect.adjustments[0] = 0.05
    rect.fill.background()
    rect.line.color.rgb = C("fail")
    rect.line.width = Pt(2)
    rect.shadow.inherit = False
    add_text(s, bx + 0.3, cy + 0.25, box_w - 0.6, 0.4,
             before.get("label", "BEFORE"), size_pt=14, weight=800, color="fail",
             allow_inline=False)
    add_text(s, bx + 0.3, cy + 0.7, box_w - 0.6, box_h - 1.0,
             before.get("quote", ""), size_pt=14, weight=500, color="ink",
             leading=1.55, allow_inline=True)

    # arrow
    add_text(s, EDGE_IN + box_w, cy, arrow_w, box_h, "→",
             size_pt=32, weight=400, color="accent",
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, allow_inline=False)

    # after box
    after = s_data.get("after", {})
    ax = EDGE_IN + box_w + arrow_w
    rect2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(ax), Inches(cy),
                               Inches(box_w), Inches(box_h))
    rect2.adjustments[0] = 0.05
    rect2.fill.background()
    rect2.line.color.rgb = C("ink_2")
    rect2.line.width = Pt(2)
    rect2.shadow.inherit = False
    add_text(s, ax + 0.3, cy + 0.25, box_w - 0.6, 0.4,
             after.get("label", "AFTER"), size_pt=14, weight=800, color="accent",
             allow_inline=False)
    add_text(s, ax + 0.3, cy + 0.7, box_w - 0.6, box_h - 1.0,
             after.get("quote", ""), size_pt=14, weight=500, color="ink",
             leading=1.55, allow_inline=True)

    # note + source
    if note:
        ny = cy + box_h + 0.3
        add_text(s, EDGE_IN, ny, avail_w, 0.8,
                 note, size_pt=13, weight=500, color="muted_2",
                 align=PP_ALIGN.CENTER, leading=1.5, allow_inline=True)
    if source:
        sy = SLIDE_H_IN - 0.95
        add_text(s, EDGE_IN, sy, avail_w, 0.3,
                 source, size_pt=10, weight=400, color="muted",
                 align=PP_ALIGN.RIGHT, allow_inline=False)
    add_footer(s, meta, page_no, total)


def slide_tagged_rows(prs, meta, s_data, page_no, total):
    """좌측 컬러 라벨 + 가로 행 비교 + 결론 (PPTX)."""
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    add_chapter_tag(s, EDGE_IN, TOP_IN, s_data.get("chapter", ""), small=True)
    add_text(s, EDGE_IN, TOP_IN + 0.4, SLIDE_W_IN - 2 * EDGE_IN, 1.2,
             s_data.get("title", ""), size_pt=36, weight=800,
             color="ink", tracking_em=-0.02, leading=1.08)
    if s_data.get("lede"):
        add_text(s, EDGE_IN, TOP_IN + 1.55, SLIDE_W_IN - 2 * EDGE_IN, 1.0,
                 s_data.get("lede", ""), size_pt=15, weight=500,
                 color="muted_2", leading=1.45, allow_inline=True)
        cy = TOP_IN + 2.4
    else:
        cy = TOP_IN + 2.0

    rows = s_data.get("rows", [])
    avail_w = SLIDE_W_IN - 2 * EDGE_IN
    tag_w = 2.4
    mid_w = (avail_w - tag_w - 0.4) * 0.55
    end_w = (avail_w - tag_w - 0.4) - mid_w
    avail_h = SLIDE_H_IN - cy - 1.7
    row_h = min(0.95, avail_h / max(1, len(rows)))
    for k, r in enumerate(rows):
        ry = cy + k * row_h
        # tag
        rect = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(EDGE_IN), Inches(ry + 0.12),
                                  Inches(tag_w), Inches(row_h - 0.24))
        rect.adjustments[0] = 0.18
        rect.fill.solid()
        rect.fill.fore_color.rgb = C("accent_soft")
        rect.line.fill.background()
        rect.shadow.inherit = False
        add_text(s, EDGE_IN, ry + 0.12, tag_w, row_h - 0.24,
                 r.get("tag", ""), size_pt=14, weight=800, color="accent",
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, allow_inline=True)
        # mid text
        add_text(s, EDGE_IN + tag_w + 0.2, ry, mid_w, row_h,
                 r.get("mid", ""), size_pt=14, weight=500, color="ink",
                 anchor=MSO_ANCHOR.MIDDLE, allow_inline=True)
        # end text
        add_text(s, EDGE_IN + tag_w + 0.2 + mid_w + 0.2, ry, end_w, row_h,
                 r.get("end", ""), size_pt=14, weight=500, color="muted_2",
                 anchor=MSO_ANCHOR.MIDDLE, allow_inline=True)
        # divider line below row
        if k < len(rows) - 1:
            add_hline(s, EDGE_IN, ry + row_h, avail_w, color="line", weight_pt=0.5)

    # conclusion
    conclusion = s_data.get("conclusion", "")
    if conclusion:
        cb_y = cy + len(rows) * row_h + 0.3
        cb_w = avail_w
        cb_x = EDGE_IN
        cb_h = 0.55
        rect = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(cb_x), Inches(cb_y),
                                  Inches(cb_w), Inches(cb_h))
        rect.adjustments[0] = 0.18
        rect.fill.background()
        rect.line.color.rgb = C("ink_2")
        rect.line.width = Pt(2)
        rect.shadow.inherit = False
        add_text(s, cb_x, cb_y, cb_w, cb_h, conclusion,
                 size_pt=13, weight=600, color="accent",
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, allow_inline=True)
    add_footer(s, meta, page_no, total)


def slide_case_profile(prs, meta, s_data, page_no, total):
    """좌 PROFILE 키-값 + 우 번호 findings + Key Insight (PPTX)."""
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    add_chapter_tag(s, EDGE_IN, TOP_IN, s_data.get("chapter", ""), small=True)
    add_text(s, EDGE_IN, TOP_IN + 0.4, SLIDE_W_IN - 2 * EDGE_IN, 1.2,
             s_data.get("title", ""), size_pt=36, weight=800,
             color="ink", tracking_em=-0.02, leading=1.08)
    if s_data.get("lede"):
        add_text(s, EDGE_IN, TOP_IN + 1.55, SLIDE_W_IN - 2 * EDGE_IN, 1.0,
                 s_data.get("lede", ""), size_pt=15, weight=500,
                 color="muted_2", leading=1.45, allow_inline=True)
        cy = TOP_IN + 2.4
    else:
        cy = TOP_IN + 2.0

    avail_w = SLIDE_W_IN - 2 * EDGE_IN
    profile_w = 4.2
    gap = 0.3
    findings_w = avail_w - profile_w - gap
    avail_h = SLIDE_H_IN - cy - 1.7
    box_h = avail_h * 0.78  # leave space for insight

    # PROFILE box (left)
    profile = s_data.get("profile", {})
    pf_rect = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(EDGE_IN), Inches(cy),
                                 Inches(profile_w), Inches(box_h))
    pf_rect.adjustments[0] = 0.04
    pf_rect.fill.solid()
    pf_rect.fill.fore_color.rgb = C("paper_2")
    pf_rect.line.fill.background()
    pf_rect.shadow.inherit = False
    pad = 0.34
    add_text(s, EDGE_IN + pad, cy + pad, profile_w - 2*pad, 0.35,
             profile.get("head", "PROFILE"), size_pt=14, weight=800, color="accent",
             tracking_em=0.18, allow_inline=False)
    add_hline(s, EDGE_IN + pad, cy + pad + 0.42, profile_w - 2*pad,
              color="line", weight_pt=0.5)
    # items
    items = profile.get("items", [])
    list_y = cy + pad + 0.55
    item_h = (box_h - (list_y - cy) - pad) / max(1, len(items))
    key_w = 0.7
    for k, it in enumerate(items):
        iy = list_y + k * item_h
        add_text(s, EDGE_IN + pad, iy, key_w, item_h,
                 it.get("key", ""), size_pt=12, weight=600, color="muted",
                 anchor=MSO_ANCHOR.MIDDLE, allow_inline=False)
        add_text(s, EDGE_IN + pad + key_w + 0.18, iy, profile_w - pad*2 - key_w - 0.18, item_h,
                 it.get("val", ""), size_pt=15, weight=700, color="ink",
                 anchor=MSO_ANCHOR.MIDDLE, allow_inline=True)

    # Findings box (right)
    fx = EDGE_IN + profile_w + gap
    f_rect = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(fx), Inches(cy),
                                Inches(findings_w), Inches(box_h))
    f_rect.adjustments[0] = 0.04
    f_rect.fill.solid()
    f_rect.fill.fore_color.rgb = C("paper")
    f_rect.line.color.rgb = C("line")
    f_rect.line.width = Pt(0.75)
    f_rect.shadow.inherit = False
    add_text(s, fx + pad, cy + pad, findings_w - 2*pad, 0.35,
             s_data.get("findings_head", "KEY POINT"), size_pt=14, weight=800, color="accent",
             tracking_em=0.18, allow_inline=False)
    add_hline(s, fx + pad, cy + pad + 0.42, findings_w - 2*pad,
              color="line", weight_pt=0.5)
    findings = s_data.get("findings", [])
    f_list_y = cy + pad + 0.55
    f_item_h = (box_h - (f_list_y - cy) - pad) / max(1, len(findings))
    for k, f in enumerate(findings):
        fy = f_list_y + k * f_item_h
        # circle number badge
        circle_d = 0.36
        oval = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                  Inches(fx + pad), Inches(fy + 0.05),
                                  Inches(circle_d), Inches(circle_d))
        oval.fill.solid()
        oval.fill.fore_color.rgb = C("ink_2")
        oval.line.fill.background()
        oval.shadow.inherit = False
        add_text(s, fx + pad, fy + 0.05, circle_d, circle_d, str(k + 1),
                 size_pt=11, weight=700, color="paper",
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, allow_inline=False)
        # title
        add_text(s, fx + pad + circle_d + 0.22, fy, findings_w - pad*2 - circle_d - 0.22, 0.4,
                 f.get("title", ""), size_pt=14, weight=700, color="ink",
                 allow_inline=True)
        # sub
        if f.get("sub"):
            add_text(s, fx + pad + circle_d + 0.22, fy + 0.42, findings_w - pad*2 - circle_d - 0.22, 0.45,
                     f.get("sub", ""), size_pt=11, weight=400, color="muted_2",
                     leading=1.5, allow_inline=True)

    # Key Insight banner
    insight = s_data.get("insight", "")
    if insight:
        in_y = cy + box_h + 0.35
        in_h = 0.6
        rect = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(EDGE_IN), Inches(in_y),
                                  Inches(avail_w), Inches(in_h))
        rect.adjustments[0] = 0.18
        rect.fill.solid()
        rect.fill.fore_color.rgb = C("highlight_yellow")
        rect.line.fill.background()
        rect.shadow.inherit = False
        add_text(s, EDGE_IN, in_y, avail_w, in_h, insight,
                 size_pt=14, weight=600, color="ink",
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, allow_inline=True)
    add_footer(s, meta, page_no, total)


def slide_concept_pill(prs, meta, s_data, page_no, total):
    """알약(pill) 컨테이너 + 원형 요소. A + B + C 형식의 '구성 공식' 슬라이드."""
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    add_chapter_tag(s, EDGE_IN, TOP_IN, s_data.get("chapter", ""), small=True)

    # Title
    add_text(s, EDGE_IN, TOP_IN + 0.4, SLIDE_W_IN - 2 * EDGE_IN, 1.2,
             s_data.get("title", ""), size_pt=36, weight=800,
             color="ink", tracking_em=-0.02, leading=1.08)
    # Lede
    has_lede = bool(s_data.get("lede"))
    if has_lede:
        add_text(s, EDGE_IN, TOP_IN + 1.55, SLIDE_W_IN - 2 * EDGE_IN, 1.0,
                 s_data.get("lede", ""), size_pt=18, weight=500,
                 color="muted_2", leading=1.5,
                 align=PP_ALIGN.CENTER, allow_inline=True)
        body_y = TOP_IN + 2.85
    else:
        body_y = TOP_IN + 2.2

    items = (s_data.get("items") or [])[:4]
    n = max(2, min(4, len(items)))
    op = s_data.get("op", "+")
    show_op = bool(op and str(op).strip())

    # 사용 가능 가로 — n별 차등 (2요소: 가장 작게 / 3요소: 컴팩트 / 4요소: 적정)
    if n == 2:
        pill_w_ratio = 0.42
    elif n == 4:
        pill_w_ratio = 0.66
    else:  # n == 3
        pill_w_ratio = 0.52
    pill_w = (SLIDE_W_IN - 2 * EDGE_IN) * pill_w_ratio
    pill_x = (SLIDE_W_IN - pill_w) / 2

    # Pill 높이 — 원 지름 + 상하 padding (전반적으로 축소)
    avail_h = SLIDE_H_IN - body_y - 1.5  # 아래 description 영역 여유
    pill_h = min(2.2, avail_h * 0.50)
    pill_pad_x = 0.4
    circle_d = pill_h - 0.6   # 원 지름 (상하 여유)
    if circle_d < 1.0:
        circle_d = 1.0

    # n개 원 + (n-1)개 + 연산자 가로 배치
    op_w = 0.4 if show_op else 0.18
    inner_w = pill_w - 2 * pill_pad_x
    circles_total = inner_w - (n - 1) * op_w
    if circles_total < n * 1.3:
        # 사이즈 부족하면 원 지름 줄이기
        circle_d = min(circle_d, circles_total / n)
    circle_w = circles_total / n  # 원 너비 (가로 점유)

    # Pill 배경 (둥근 사각형)
    pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(pill_x), Inches(body_y),
                              Inches(pill_w), Inches(pill_h))
    pill.adjustments[0] = 0.5  # 가장 둥글게 (= stadium)
    pill.fill.solid()
    pill.fill.fore_color.rgb = C("ink_2")
    pill.line.fill.background()
    pill.shadow.inherit = False

    # Circles + operators
    cx = pill_x + pill_pad_x
    cy = body_y + (pill_h - circle_d) / 2
    desc_y = body_y + pill_h + 0.25
    desc_h = 1.0
    for j, it in enumerate(items):
        # circle x: centered within its width slot
        slot_x = cx + j * (circle_w + op_w)
        circle_x = slot_x + (circle_w - circle_d) / 2
        oval = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                  Inches(circle_x), Inches(cy),
                                  Inches(circle_d), Inches(circle_d))
        oval.fill.solid()
        oval.fill.fore_color.rgb = C("paper_2")
        oval.line.fill.background()
        oval.shadow.inherit = False
        # circle text
        font_size = 18 if n == 3 else (16 if n == 4 else 22)
        add_text(s, circle_x + 0.1, cy + 0.1,
                 circle_d - 0.2, circle_d - 0.2,
                 it.get("circle", ""), size_pt=font_size, weight=800,
                 color="ink", leading=1.2, tracking_em=-0.02,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 allow_inline=True)
        # operator (between circles) — op이 비어있으면 생략
        if j < n - 1 and show_op:
            op_x = slot_x + circle_w
            add_text(s, op_x, cy, op_w, circle_d,
                     op, size_pt=20, weight=700, color="paper",
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                     allow_inline=False)
        # description below circle (서브 텍스트 키움 14 → 17)
        add_text(s, slot_x, desc_y, circle_w, desc_h,
                 it.get("desc", ""), size_pt=17, weight=500,
                 color="muted_2", leading=1.45,
                 align=PP_ALIGN.CENTER, allow_inline=True)

    add_footer(s, meta, page_no, total)


def slide_alert_close(prs, meta, s_data, page_no, total):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    add_chapter_tag(s, EDGE_IN, TOP_IN, s_data.get("chapter", ""), small=True)
    add_text(s, EDGE_IN, TOP_IN + 0.4, SLIDE_W_IN - 2 * EDGE_IN, 1.6,
             s_data.get("title", ""), size_pt=40, weight=800,
             color="ink", tracking_em=-0.02, leading=1.08)
    add_text(s, EDGE_IN, TOP_IN + 1.7, SLIDE_W_IN - 2 * EDGE_IN, 0.9,
             s_data.get("lede", ""), size_pt=18, weight=500,
             color="muted_2", leading=1.45, allow_inline=False)
    # 슬라이드 가운데보다 살짝 위에 큰 알림 메시지 (2줄 줄바꿈 자동 처리)
    body_y = TOP_IN + 2.7
    body_h = SLIDE_H_IN - body_y - 1.85  # 아래 padding 늘려 위쪽으로 정렬되게
    add_text(s, EDGE_IN, body_y, SLIDE_W_IN - 2 * EDGE_IN, body_h,
             s_data.get("alert", ""), size_pt=42, weight=700,
             color="alert_red", leading=1.4, tracking_em=-0.02,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_footer(s, meta, page_no, total)


# ============================================================
#  Design system meta — color_palette / type_scale (PPTX)
# ============================================================

def slide_color_palette(prs, meta, s_data, page_no, total):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    add_chapter_tag(s, EDGE_IN, TOP_IN, s_data.get("chapter", ""), small=True)
    add_text(s, EDGE_IN, TOP_IN + 0.4, SLIDE_W_IN - 2 * EDGE_IN, 1.2,
             s_data.get("title", ""), size_pt=36, weight=800,
             color="ink", tracking_em=-0.02, leading=1.08)
    if s_data.get("lede"):
        add_text(s, EDGE_IN, TOP_IN + 1.6, SLIDE_W_IN - 2 * EDGE_IN, 0.6,
                 s_data["lede"], size_pt=15, weight=500, color="muted_2",
                 leading=1.4, allow_inline=False)
        cy = TOP_IN + 2.555
    else:
        cy = TOP_IN + 2.2
    avail_w = SLIDE_W_IN - 2 * EDGE_IN
    avail_h = SLIDE_H_IN - cy - 0.85
    swatches = s_data.get("swatches", [])
    n = max(3, min(5, s_data.get("cols", 4)))
    inner_gap = 0.16
    sw_w = (avail_w - inner_gap * (n - 1)) / n
    sw_h = avail_h
    chip_h = sw_h * 0.55
    for j, sw in enumerate(swatches[:n]):
        x = EDGE_IN + j * (sw_w + inner_gap)
        # outer frame
        outer = add_rect(s, x, cy, sw_w, sw_h, fill="paper", line="line", line_w_pt=0.5)
        # color chip
        hex_val = sw.get("hex", "#000000")
        rgb = hex_to_rgb(hex_val)
        chip = add_rect(s, x, cy, sw_w, chip_h, fill="paper")
        chip.fill.fore_color.rgb = rgb
        chip.line.fill.background()
        # name on chip (bottom-left of chip)
        on_color = sw.get("on", "#FFFFFF")
        rPr_color = hex_to_rgb(on_color)
        # use a textbox with custom color
        tb = s.shapes.add_textbox(Inches(x + 0.18), Inches(cy + chip_h - 0.5),
                                  Inches(sw_w - 0.36), Inches(0.4))
        tf = tb.text_frame
        tf.margin_left = tf.margin_right = 0
        tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        for old in list(p.runs):
            old._r.getparent().remove(old._r)
        r = p.add_run()
        r.text = sw.get("name", "")
        _set_font_typeface(r, FONT_FAMILY)
        r.font.size = Pt(11)
        r.font.bold = True
        r.font.color.rgb = rPr_color
        # body info area
        body_y = cy + chip_h
        body_h = sw_h - chip_h
        add_hline(s, x, body_y, sw_w, color="line", weight_pt=0.5)
        add_text(s, x + 0.18, body_y + 0.15, sw_w - 0.36, 0.35,
                 sw.get("role", ""), size_pt=10, weight=600, color="ink",
                 leading=1.3, allow_inline=False)
        add_text(s, x + 0.18, body_y + 0.45, sw_w - 0.36, 0.3,
                 hex_val, size_pt=9, weight=500, color="muted_2",
                 leading=1.3, allow_inline=False)
        if sw.get("token"):
            add_text(s, x + 0.18, body_y + 0.7, sw_w - 0.36, 0.3,
                     sw["token"], size_pt=9, weight=500, color="muted",
                     leading=1.3, allow_inline=False)
    add_footer(s, meta, page_no, total)


def slide_type_scale(prs, meta, s_data, page_no, total):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    add_chapter_tag(s, EDGE_IN, TOP_IN, s_data.get("chapter", ""), small=True)
    add_text(s, EDGE_IN, TOP_IN + 0.4, SLIDE_W_IN - 2 * EDGE_IN, 1.2,
             s_data.get("title", ""), size_pt=36, weight=800,
             color="ink", tracking_em=-0.02, leading=1.08)
    if s_data.get("lede"):
        add_text(s, EDGE_IN, TOP_IN + 1.6, SLIDE_W_IN - 2 * EDGE_IN, 0.6,
                 s_data["lede"], size_pt=15, weight=500, color="muted_2",
                 leading=1.4, allow_inline=False)
        cy = TOP_IN + 2.55
    else:
        cy = TOP_IN + 2.1
    avail_w = SLIDE_W_IN - 2 * EDGE_IN
    avail_h = SLIDE_H_IN - cy - 0.85
    items = s_data.get("items", [])
    n = max(1, len(items))
    row_h = min(0.85, avail_h / n)
    meta_w = avail_w * 0.32
    sample_w = avail_w - meta_w
    # top border
    add_hline(s, EDGE_IN, cy, avail_w, color="line", weight_pt=0.5)
    for j, it in enumerate(items):
        ry = cy + j * row_h
        # row bg
        bg = add_rect(s, EDGE_IN, ry, avail_w, row_h, fill="paper")
        bg.line.fill.background()
        # bottom divider
        add_hline(s, EDGE_IN, ry + row_h, avail_w, color="line", weight_pt=0.5)
        # meta
        name = it.get("name", "")
        size_px = it.get("size_px", 18)
        weight = it.get("weight", 400)
        leading = it.get("leading", 1.4)
        tracking = it.get("tracking_em", 0)
        add_text(s, EDGE_IN + 0.2, ry + 0.1, meta_w - 0.4, 0.4,
                 name, size_pt=12, weight=700, color="ink",
                 tracking_em=-0.01, allow_inline=False)
        meta_str = f"{size_px}px · {weight} · lh {leading} · tr {tracking}em"
        add_text(s, EDGE_IN + 0.2, ry + 0.45, meta_w - 0.4, 0.3,
                 meta_str, size_pt=9, weight=400, color="muted_2",
                 allow_inline=False)
        # sample (using PPT pt scale: PPTX 캔버스에서 px → pt 변환 ≈ size_px * 0.55)
        sample = it.get("sample", "샘플 텍스트")
        sample_pt = int(size_px * 0.55)
        sample_pt = min(48, max(10, sample_pt))
        color_key = it.get("color_key", "ink")
        add_text(s, EDGE_IN + meta_w + 0.3, ry + 0.05, sample_w - 0.4, row_h - 0.1,
                 sample, size_pt=sample_pt, weight=weight, color=color_key,
                 tracking_em=tracking, leading=leading,
                 anchor=MSO_ANCHOR.MIDDLE)
    add_footer(s, meta, page_no, total)


# ============================================================
#  6 new templates (PPTX)
# ============================================================

def slide_checklist(prs, meta, s_data, page_no, total):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    add_chapter_tag(s, EDGE_IN, TOP_IN, s_data.get("chapter", ""), small=True)
    add_text(s, EDGE_IN, TOP_IN + 0.4, SLIDE_W_IN - 2 * EDGE_IN, 1.4,
             s_data.get("title", ""), size_pt=36, weight=800,
             color="ink", tracking_em=-0.02, leading=1.08)
    cy = TOP_IN + 2.2
    box_h = SLIDE_H_IN - cy - 0.85
    box_w = SLIDE_W_IN - 2 * EDGE_IN
    # outer frame
    frame = add_rect(s, EDGE_IN, cy, box_w, box_h, fill="paper", line="ink", line_w_pt=1.0)
    pad_x, pad_y = 0.5, 0.5
    inner_x = EDGE_IN + pad_x
    inner_y = cy + pad_y
    inner_w = box_w - 2 * pad_x
    # Box title
    if s_data.get("box_title"):
        add_text(s, inner_x, inner_y, inner_w, 0.7,
                 s_data["box_title"], size_pt=24, weight=800, color="ink",
                 tracking_em=-0.01, leading=1.2)
        inner_y += 0.7
    if s_data.get("subtitle"):
        add_text(s, inner_x, inner_y, inner_w, 0.55,
                 s_data["subtitle"], size_pt=15, weight=400, color="ink",
                 leading=1.4, allow_inline=False)
        inner_y += 0.7
    inner_y += 0.2
    items = s_data.get("items", [])
    item_h = 0.5
    for it in items:
        if isinstance(it, dict):
            text = it.get("text", "")
            emphasis = it.get("emphasis", False)
        else:
            text = str(it); emphasis = False
        weight = 700 if emphasis else 400
        add_text(s, inner_x, inner_y, inner_w, item_h,
                 text, size_pt=16, weight=weight, color="ink", leading=1.5)
        if emphasis:
            # 밑줄 효과: 텍스트 아래 얇은 선 (대략적)
            add_hline(s, inner_x, inner_y + 0.3, min(6.5, len(text) * 0.13),
                      color="ink", weight_pt=1.0)
        inner_y += 0.45
    add_footer(s, meta, page_no, total)


def slide_summary_takeaway(prs, meta, s_data, page_no, total):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    # 표준 head 포맷: chapter-tag (small) + slide-title + 옵션 lede
    chapter_text = s_data.get("chapter") or s_data.get("section_label", "")
    add_chapter_tag(s, EDGE_IN, TOP_IN, chapter_text, small=True)
    add_text(s, EDGE_IN, TOP_IN + 0.4, SLIDE_W_IN - 2 * EDGE_IN, 1.4,
             s_data.get("title", ""), size_pt=36, weight=800,
             color="ink", tracking_em=-0.02, leading=1.08)
    if s_data.get("lede"):
        add_text(s, EDGE_IN, TOP_IN + 1.6, SLIDE_W_IN - 2 * EDGE_IN, 0.6,
                 s_data["lede"], size_pt=15, weight=500, color="muted_2",
                 leading=1.4, allow_inline=False)
        cy = TOP_IN + 2.55
    else:
        cy = TOP_IN + 2.1
    cards = s_data.get("cards", [])[:4]
    n = max(2, min(4, len(cards)))
    avail_w = SLIDE_W_IN - 2 * EDGE_IN
    inner_gap = 0.22
    card_w = (avail_w - inner_gap * (n - 1)) / n
    # Reserve takeaway bar height
    takeaway_h = 0.7
    bottom_pad = 0.85
    cards_h = SLIDE_H_IN - cy - takeaway_h - bottom_pad - 0.3
    for j, c in enumerate(cards):
        cx = EDGE_IN + j * (card_w + inner_gap)
        # Light pastel header band
        head_h = 0.85
        head = add_rect(s, cx, cy, card_w, head_h, fill="pastel_blue")
        head.line.fill.background()
        add_text(s, cx + 0.15, cy, card_w - 0.3, head_h,
                 c.get("title", ""), size_pt=18, weight=700, color="ink",
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 leading=1.25, allow_inline=False)
        # Body items
        iy = cy + head_h + 0.3
        for it in c.get("items", [])[:6]:
            add_text(s, cx, iy, 0.18, 0.4, "•",
                     size_pt=14, weight=700, color="accent", allow_inline=False)
            add_text(s, cx + 0.2, iy, card_w - 0.2, 0.6,
                     it, size_pt=13, weight=400, color="ink", leading=1.5)
            iy += 0.5
        # Bottom border for card
        add_hline(s, cx, cy + cards_h, card_w, color="ink", weight_pt=0.75)
    # Takeaway bar
    tk_y = SLIDE_H_IN - bottom_pad - takeaway_h
    bar = add_rect(s, EDGE_IN, tk_y, avail_w, takeaway_h, fill="accent")
    add_text(s, EDGE_IN, tk_y, avail_w, takeaway_h,
             s_data.get("takeaway", ""), size_pt=18, weight=700, color="paper",
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             tracking_em=-0.01)
    add_footer(s, meta, page_no, total)


def slide_dual_panel(prs, meta, s_data, page_no, total):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    # 표준 head 포맷: chapter-tag + slide-title + 옵션 lede
    chapter_text = s_data.get("chapter") or s_data.get("section_label", "")
    add_chapter_tag(s, EDGE_IN, TOP_IN, chapter_text, small=True)
    add_text(s, EDGE_IN, TOP_IN + 0.4, SLIDE_W_IN - 2 * EDGE_IN, 1.4,
             s_data.get("title", ""), size_pt=36, weight=800,
             color="ink", tracking_em=-0.02, leading=1.08)
    if s_data.get("lede"):
        add_text(s, EDGE_IN, TOP_IN + 1.6, SLIDE_W_IN - 2 * EDGE_IN, 0.6,
                 s_data["lede"], size_pt=15, weight=500, color="muted_2",
                 leading=1.4, allow_inline=False)
        cy = TOP_IN + 2.55
    else:
        cy = TOP_IN + 2.1
    pane_h = SLIDE_H_IN - cy - 0.85
    pane_w = (SLIDE_W_IN - 2 * EDGE_IN) / 2

    def _pane(x, color_key, pane_data, tint):
        head_h = 0.55
        head = add_rect(s, x, cy, pane_w, head_h, fill=color_key)
        head.line.fill.background()
        add_text(s, x, cy, pane_w, head_h,
                 pane_data.get("header", ""),
                 size_pt=18, weight=800, color="paper",
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 tracking_em=-0.01, allow_inline=False)
        # body bg (very subtle tint)
        body_y = cy + head_h
        body_h = pane_h - head_h
        body = add_rect(s, x, body_y, pane_w, body_h, fill=tint)
        body.line.color.rgb = C("line")
        body.line.width = Pt(0.5)
        # sections
        sec_x = x + 0.3
        sec_y = body_y + 0.25
        sec_w = pane_w - 0.6
        sections = pane_data.get("sections", [])
        for si, sec in enumerate(sections):
            if si > 0:
                add_hline(s, sec_x, sec_y, sec_w, color="line", weight_pt=0.75)
                sec_y += 0.18
            add_text(s, sec_x, sec_y, sec_w, 0.55,
                     sec.get("title", ""), size_pt=16, weight=700,
                     color=color_key, leading=1.2)
            sec_y += 0.5
            for it in sec.get("items", [])[:6]:
                add_text(s, sec_x, sec_y, 0.18, 0.35, "•",
                         size_pt=13, weight=700, color="ink", allow_inline=False)
                add_text(s, sec_x + 0.2, sec_y, sec_w - 0.2, 0.6,
                         it, size_pt=13, weight=400, color="ink", leading=1.45)
                sec_y += 0.4
            sec_y += 0.15
        # footnote
        if pane_data.get("footnote"):
            add_text(s, sec_x, body_y + body_h - 0.6, sec_w, 0.5,
                     "* " + pane_data["footnote"],
                     size_pt=10, weight=500, color="muted_2", leading=1.4,
                     allow_inline=False)

    _pane(EDGE_IN,            "accent", s_data.get("left", {}),  "paper_2")
    _pane(EDGE_IN + pane_w,   "alert_red", s_data.get("right", {}), "pastel_red")
    add_footer(s, meta, page_no, total)


def slide_case_grid_4(prs, meta, s_data, page_no, total):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    chapter_text = s_data.get("chapter") or s_data.get("section_label", "")
    add_chapter_tag(s, EDGE_IN, TOP_IN, chapter_text, small=True)
    add_text(s, EDGE_IN, TOP_IN + 0.4, SLIDE_W_IN - 2 * EDGE_IN, 1.4,
             s_data.get("title", ""), size_pt=36, weight=800, color="ink",
             tracking_em=-0.02, leading=1.08)
    if s_data.get("lede"):
        add_text(s, EDGE_IN, TOP_IN + 1.6, SLIDE_W_IN - 2 * EDGE_IN, 0.6,
                 s_data["lede"], size_pt=15, weight=500, color="muted_2",
                 leading=1.4, allow_inline=False)
        cy = TOP_IN + 2.555
    else:
        cy = TOP_IN + 2.2
    avail_w = SLIDE_W_IN - 2 * EDGE_IN
    avail_h = SLIDE_H_IN - cy - 0.85
    cw = avail_w / 2
    ch = avail_h / 2
    cards = s_data.get("cards", [])[:4]
    while len(cards) < 4:
        cards.append({})
    # Outer top border
    add_hline(s, EDGE_IN, cy, avail_w, color="ink", weight_pt=1.0)
    for j, c in enumerate(cards):
        row, col = divmod(j, 2)
        x = EDGE_IN + col * cw
        y = cy + row * ch
        odd_idx = (j % 2 == 0)  # 좌상=0, 우상=1, 좌하=2, 우하=3 → odd by index 0,2 → blue
        is_blue = (j == 0 or j == 3)
        fill = "pastel_blue" if is_blue else "pastel_red"
        title_color = "accent" if is_blue else "alert_red"
        # cell bg
        bg = add_rect(s, x, y, cw, ch, fill=fill)
        bg.line.fill.background()
        # bottom border
        add_hline(s, x, y + ch, cw, color="ink", weight_pt=1.0)
        # right border (only for left cells)
        if col == 0:
            ln = s.shapes.add_connector(1, Inches(x + cw), Inches(y),
                                        Inches(x + cw), Inches(y + ch))
            ln.line.color.rgb = C("ink"); ln.line.width = Pt(1.0)
        # content
        ix = x + 0.32; iy = y + 0.3; iw = cw - 0.64
        add_text(s, ix, iy, iw, 0.65,
                 c.get("title", ""), size_pt=18, weight=700, color=title_color,
                 leading=1.25)
        iy += 0.6
        for it in c.get("items", [])[:5]:
            add_text(s, ix, iy, 0.18, 0.35, "•",
                     size_pt=12, weight=700, color="ink", allow_inline=False)
            add_text(s, ix + 0.2, iy, iw - 0.2, 0.5,
                     it, size_pt=12, weight=400, color="ink", leading=1.5)
            iy += 0.38
        if c.get("arrow"):
            iy += 0.05
            add_hline(s, ix, iy, iw - 0.1, color="line", weight_pt=0.5)
            iy += 0.08
            add_text(s, ix, iy, iw, 0.5,
                     "→ " + c["arrow"],
                     size_pt=13, weight=700, color=title_color, leading=1.4)
    add_footer(s, meta, page_no, total)


def slide_case_analysis(prs, meta, s_data, page_no, total):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    add_chapter_tag(s, EDGE_IN, TOP_IN, s_data.get("chapter", ""), small=True)
    add_text(s, EDGE_IN, TOP_IN + 0.4, SLIDE_W_IN - 2 * EDGE_IN, 1.2,
             s_data.get("title", ""), size_pt=36, weight=800,
             color="ink", tracking_em=-0.02, leading=1.08)
    cy = TOP_IN + 1.9
    avail_w = SLIDE_W_IN - 2 * EDGE_IN
    avail_h = SLIDE_H_IN - cy - 1.0   # footer 영역 더 여유롭게 (0.85→1.0)
    gap = 0.4
    ref_w = avail_w * 0.5             # 우측 영역 더 넓게 (0.55→0.5)
    pts_x = EDGE_IN + ref_w + gap
    pts_w = avail_w - ref_w - gap
    # Reference panel
    ref_bg = add_rect(s, EDGE_IN, cy, ref_w, avail_h, fill="paper_2")
    ref_bg.line.fill.background()
    if s_data.get("ref_title"):
        add_text(s, EDGE_IN + 0.3, cy + 0.2, ref_w - 0.6, 0.5,
                 s_data["ref_title"], size_pt=14, weight=700, color="ink",
                 tracking_em=-0.01)
    if s_data.get("ref_image"):
        img_p = _resolve_image(s_data["ref_image"])
        _add_picture_fit(s, img_p, EDGE_IN + 0.3, cy + 0.85,
                         ref_w - 0.6, avail_h - 1.05)
    elif s_data.get("ref_text"):
        # 글자 더 작게 + auto-shrink 효과 위해 size 9로
        add_text(s, EDGE_IN + 0.3, cy + 0.85, ref_w - 0.6, avail_h - 1.05,
                 s_data["ref_text"], size_pt=9, weight=400, color="ink",
                 leading=1.45, allow_inline=False)
    # Findings panel — conclusion 영역을 충분히 확보 (1.2 inch)
    findings = s_data.get("findings", [])
    iy = cy + 0.1
    concl_h = 1.2 if s_data.get("conclusion") else 0   # 0.7 → 1.2 (텍스트 2줄 여유)
    avail_for_findings = avail_h - concl_h - 0.3       # 마진 0.15 → 0.3
    f_h = avail_for_findings / max(1, len(findings))
    for j, f in enumerate(findings):
        label = f.get("label", f"합격 포인트 {j+1}")
        label_w = max(1.1, len(label) * 0.13 + 0.4)
        bg_label = add_rect(s, pts_x, iy, label_w, 0.35, fill="accent")
        bg_label.line.fill.background()
        add_text(s, pts_x, iy, label_w, 0.35, label,
                 size_pt=11, weight=700, color="paper",
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 tracking_em=-0.01, allow_inline=False)
        add_text(s, pts_x, iy + 0.45, pts_w, f_h - 0.5,
                 f.get("text", ""), size_pt=14, weight=500, color="ink",
                 leading=1.5)
        iy += f_h
    # Conclusion: 위쪽 강조 라인 + 네이비 굵은 텍스트
    if s_data.get("conclusion"):
        cb_y = cy + avail_h - concl_h
        ln = s.shapes.add_connector(1, Inches(pts_x), Inches(cb_y),
                                    Inches(pts_x + pts_w), Inches(cb_y))
        ln.line.color.rgb = C("accent")
        ln.line.width = Pt(2.0)
        # 텍스트 size 18 + height 1.0 → 2줄까지 안전하게 들어감
        add_text(s, pts_x, cb_y + 0.15, pts_w, concl_h - 0.2,
                 s_data["conclusion"], size_pt=18, weight=700, color="accent",
                 anchor=MSO_ANCHOR.TOP, leading=1.4, tracking_em=-0.01)
    add_footer(s, meta, page_no, total)


def slide_pipeline_matrix(prs, meta, s_data, page_no, total):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    add_chapter_tag(s, EDGE_IN, TOP_IN, s_data.get("chapter", ""), small=True)
    add_text(s, EDGE_IN, TOP_IN + 0.4, SLIDE_W_IN - 2 * EDGE_IN, 1.4,
             s_data.get("title", ""), size_pt=32, weight=800,
             color="ink", tracking_em=-0.02, leading=1.08)
    if s_data.get("lede"):
        add_text(s, EDGE_IN, TOP_IN + 1.6, SLIDE_W_IN - 2 * EDGE_IN, 0.6,
                 s_data["lede"], size_pt=15, weight=500, color="muted_2",
                 leading=1.4, allow_inline=False)
        cy = TOP_IN + 2.555
    else:
        cy = TOP_IN + 2.1
    stages = s_data.get("stages", [])[:3]
    while len(stages) < 3:
        stages.append({"header": "", "rows": []})
    # Compute layout
    avail_w = SLIDE_W_IN - 2 * EDGE_IN
    arrow_w = 0.45
    # Reserve note area on right side if needed
    if s_data.get("note"):
        avail_w_pipeline = avail_w - 1.6
    else:
        avail_w_pipeline = avail_w
    stage_w = (avail_w_pipeline - 2 * arrow_w) / 3
    avail_h = SLIDE_H_IN - cy - 0.85
    head_h = 0.7
    body_h = avail_h - head_h - 0.2
    # Header row
    for si in range(3):
        sx = EDGE_IN + si * (stage_w + arrow_w)
        bg = add_rect(s, sx, cy, stage_w, head_h, fill="ink")
        bg.line.fill.background()
        add_text(s, sx + 0.1, cy, stage_w - 0.2, head_h,
                 stages[si].get("header", ""),
                 size_pt=15, weight=700, color="paper",
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 tracking_em=-0.01, leading=1.25)
        if si < 2:
            arrow_x = sx + stage_w
            add_text(s, arrow_x, cy, arrow_w, head_h,
                     "→", size_pt=22, weight=700, color="ink",
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                     allow_inline=False)
    # Body rows
    max_rows = max((len(stages[si].get("rows", [])) for si in range(3)), default=0)
    if max_rows > 0:
        body_y = cy + head_h + 0.15
        row_h = body_h / max_rows
        for ri in range(max_rows):
            ry = body_y + ri * row_h
            for si in range(3):
                sx = EDGE_IN + si * (stage_w + arrow_w)
                rows = stages[si].get("rows", [])
                cell = rows[ri] if ri < len(rows) else None
                if cell is None:
                    continue
                cell_pad = 0.05
                inner_h = row_h - 2 * cell_pad
                if isinstance(cell, str):
                    bg = add_rect(s, sx, ry + cell_pad, stage_w, inner_h, fill="ink")
                    bg.line.fill.background()
                    add_text(s, sx + 0.1, ry + cell_pad, stage_w - 0.2, inner_h,
                             cell, size_pt=14, weight=700, color="paper",
                             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                             leading=1.3)
                elif isinstance(cell, list):
                    n = len(cell)
                    sub_gap = 0.05
                    sub_w = (stage_w - sub_gap * (n - 1)) / n
                    for ci, c in enumerate(cell):
                        cx = sx + ci * (sub_w + sub_gap)
                        bg = add_rect(s, cx, ry + cell_pad, sub_w, inner_h, fill="ink")
                        bg.line.fill.background()
                        label = c if isinstance(c, str) else c.get("label", "")
                        add_text(s, cx + 0.05, ry + cell_pad, sub_w - 0.1, inner_h,
                                 label, size_pt=12, weight=700, color="paper",
                                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                                 leading=1.25)
                elif isinstance(cell, dict):
                    label = cell.get("label", "")
                    sub = cell.get("sub", [])
                    if sub:
                        # main label on top, sub-cells below
                        top_h = inner_h * 0.45
                        bot_h = inner_h - top_h - 0.04
                        bg = add_rect(s, sx, ry + cell_pad, stage_w, top_h, fill="ink")
                        bg.line.fill.background()
                        add_text(s, sx + 0.1, ry + cell_pad, stage_w - 0.2, top_h,
                                 label, size_pt=13, weight=700, color="paper",
                                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                                 leading=1.25)
                        n = len(sub); sub_gap = 0.05
                        sub_w = (stage_w - sub_gap * (n - 1)) / n
                        for ci, c in enumerate(sub):
                            cx = sx + ci * (sub_w + sub_gap)
                            cy_s = ry + cell_pad + top_h + 0.04
                            bg2 = add_rect(s, cx, cy_s, sub_w, bot_h, fill="ink")
                            bg2.line.fill.background()
                            slabel = c if isinstance(c, str) else c.get("label", "")
                            add_text(s, cx + 0.04, cy_s, sub_w - 0.08, bot_h,
                                     slabel, size_pt=10, weight=600, color="paper",
                                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                                     leading=1.25)
                    else:
                        bg = add_rect(s, sx, ry + cell_pad, stage_w, inner_h, fill="ink")
                        bg.line.fill.background()
                        add_text(s, sx + 0.1, ry + cell_pad, stage_w - 0.2, inner_h,
                                 label, size_pt=14, weight=700, color="paper",
                                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                                 leading=1.3)
    # Side note
    if s_data.get("note"):
        nx = EDGE_IN + avail_w_pipeline + 0.3
        nw = avail_w - avail_w_pipeline - 0.3
        add_text(s, nx, cy + 0.2, nw, 1.4,
                 s_data["note"], size_pt=14, weight=700, color="ink",
                 leading=1.4, align=PP_ALIGN.CENTER, allow_inline=False)
    add_footer(s, meta, page_no, total)


BUILDERS = {
    "cover":           lambda prs, meta, s, p, t: slide_cover(prs, meta),
    "chapter_divider": slide_chapter_divider,
    "content":         slide_content,
    "card_grid_4":     slide_card_grid_4,
    "card_grid":       slide_card_grid,
    "two_col":         slide_two_col,
    "toc":             slide_toc,
    "speaker":         slide_speaker,
    "stage_flow":      slide_stage_flow,
    "alert_close":     slide_alert_close,
    "qa_close":        slide_qa_close,
    "compare_table":   slide_compare_table,
    "image_full":      slide_image_full,
    "image_headed":    slide_image_headed,
    "image_split":     slide_image_split,
    "checklist":        slide_checklist,
    "summary_takeaway": slide_summary_takeaway,
    "dual_panel":       slide_dual_panel,
    "case_grid_4":      slide_case_grid_4,
    "case_analysis":    slide_case_analysis,
    "pipeline_matrix":  slide_pipeline_matrix,
    "color_palette":    slide_color_palette,
    "type_scale":       slide_type_scale,
    "concept_pill":     slide_concept_pill,
    "priority_matrix":  slide_priority_matrix,
    "step_compare":     slide_step_compare,
    "before_after":     slide_before_after,
    "tagged_rows":      slide_tagged_rows,
    "case_profile":     slide_case_profile,
}


def build(data: dict) -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    meta = data.get("meta", {})
    slides = data.get("slides", [])
    total = len(slides)
    for i, s in enumerate(slides, 1):
        kind = s.get("type")
        fn = BUILDERS.get(kind)
        if not fn:
            print(f"[warn] unknown slide type: {kind}", file=sys.stderr)
            continue
        fn(prs, meta, s, i, total)
    return prs


def main():
    if len(sys.argv) < 2:
        print("usage: build_pptx.py <input.json> [output.pptx]", file=sys.stderr)
        sys.exit(2)
    in_path = Path(sys.argv[1])
    out_path = Path(sys.argv[2]) if len(sys.argv) >= 3 else OUTPUTS / (in_path.stem + ".pptx")
    out_path.parent.mkdir(parents=True, exist_ok=True)
    data = json.loads(in_path.read_text(encoding="utf-8"))
    prs = build(data)
    prs.save(out_path)
    print(f"OK  PPTX  -> {out_path}")


if __name__ == "__main__":
    main()
